import os
import re
import fitz  # PyMuPDF
from pptx import Presentation

# Try to import OCR support (optional)
try:
    from pymupdf4llm import to_markdown
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    to_markdown = None

def normalize_text(text: str) -> str:
    """
    Normalize extracted text:
    - Remove excessive whitespace
    - Fix line breaks in English text
    - Remove special characters that break text flow
    - Preserve structure markers
    """
    if not text:
        return ""
    
    # Replace multiple spaces with single space
    text = re.sub(r' +', ' ', text)
    
    # Fix common PDF extraction issues:
    # Remove spaces before punctuation (except for abbreviations)
    text = re.sub(r'\s+([.,!?;:])', r'\1', text)
    
    # Fix line breaks in the middle of words (common in PDF extraction)
    # Remove single newlines that break words, but keep double newlines
    lines = text.split('\n')
    normalized_lines = []
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            normalized_lines.append('')
            continue
        
        # If line doesn't end with punctuation and next line starts with lowercase,
        # it's likely a broken word - join them
        if (i < len(lines) - 1 and 
            line and 
            not line[-1] in '.!?;:' and
            lines[i + 1].strip() and
            lines[i + 1].strip()[0].islower()):
            # Join with space instead of newline
            normalized_lines.append(line + ' ')
        else:
            normalized_lines.append(line)
    
    text = '\n'.join(normalized_lines)
    
    # Remove excessive newlines (more than 2 consecutive)
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Remove trailing whitespace from each line
    lines = [line.rstrip() for line in text.split('\n')]
    text = '\n'.join(lines)
    
    return text.strip()


def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from PDF or PPTX file.
    
    Args:
        file_path: Path to the file
        
    Returns:
        Extracted and normalized text as string
        
    Raises:
        ValueError: If file format is not supported or file doesn't exist
    """
    # Check if file exists
    if not os.path.exists(file_path):
        raise ValueError(f"Файл не найден: {file_path}")
    
    # Check file size
    file_size = os.path.getsize(file_path)
    if file_size == 0:
        raise ValueError("Файл пустой")
    
    if file_path.lower().endswith(".pdf"):
        raw_text = extract_text_from_pdf(file_path)
    elif file_path.lower().endswith(".pptx"):
        raw_text = extract_text_from_pptx(file_path)
    else:
        raise ValueError("Поддерживаются только PDF и PPTX файлы.")
    
    # Normalize the extracted text
    normalized_text = normalize_text(raw_text)
    
    return normalized_text


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file using PyMuPDF with multiple extraction methods."""
    text = ""
    total_pages = 0
    pages_with_text = 0
    
    try:
        with fitz.open(file_path) as doc:
            total_pages = len(doc)
            
            for page_num, page in enumerate(doc, 1):
                page_text = ""
                
                # Method 1: Standard text extraction
                try:
                    page_text = page.get_text("text")
                except Exception:
                    pass
                
                # Method 2: Try blocks extraction if standard method returned little text
                if len(page_text.strip()) < 50:
                    try:
                        blocks = page.get_text("blocks")
                        if blocks:
                            block_text = "\n".join([block[4] for block in blocks if block[4].strip()])
                            if len(block_text) > len(page_text):
                                page_text = block_text
                    except Exception:
                        pass
                
                # Method 3: Try dict extraction for structured text
                if len(page_text.strip()) < 50:
                    try:
                        text_dict = page.get_text("dict")
                        if text_dict and "blocks" in text_dict:
                            dict_text_parts = []
                            for block in text_dict["blocks"]:
                                if block.get("type") == 0:  # Text block
                                    if "lines" in block:
                                        for line in block["lines"]:
                                            if "spans" in line:
                                                line_text = " ".join([
                                                    span.get("text", "") 
                                                    for span in line["spans"]
                                                    if span.get("text", "").strip()
                                                ])
                                                if line_text.strip():
                                                    dict_text_parts.append(line_text)
                            if dict_text_parts:
                                dict_text = "\n".join(dict_text_parts)
                                if len(dict_text) > len(page_text):
                                    page_text = dict_text
                    except Exception:
                        pass
                
                # Method 4: Try rawdict extraction (most detailed)
                if len(page_text.strip()) < 50:
                    try:
                        raw_dict = page.get_text("rawdict")
                        if raw_dict and "blocks" in raw_dict:
                            raw_text_parts = []
                            for block in raw_dict["blocks"]:
                                if block.get("type") == 0:  # Text block
                                    if "lines" in block:
                                        for line in block["lines"]:
                                            if "spans" in line:
                                                for span in line["spans"]:
                                                    span_text = span.get("text", "").strip()
                                                    if span_text:
                                                        raw_text_parts.append(span_text)
                            if raw_text_parts:
                                raw_text = " ".join(raw_text_parts)
                                if len(raw_text) > len(page_text):
                                    page_text = raw_text
                    except Exception:
                        pass
                
                # Method 5: Try to extract from annotations (comments, notes)
                if len(page_text.strip()) < 50:
                    try:
                        annot_texts = []
                        for annot in page.annots():
                            if annot.info and annot.info.get("content"):
                                annot_texts.append(annot.info["content"])
                        if annot_texts:
                            annot_text = "\n".join(annot_texts)
                            if len(annot_text) > len(page_text):
                                page_text = annot_text
                    except Exception:
                        pass
                
                # Clean up the page text
                if page_text.strip():
                    # Remove page numbers and headers/footers (common patterns)
                    page_text = re.sub(r'^\s*\d+\s*$', '', page_text, flags=re.MULTILINE)
                    
                    text += f"\n--- Страница {page_num} ---\n"
                    text += page_text + "\n"
                    pages_with_text += 1
                    
    except Exception as e:
        raise ValueError(f"Ошибка при чтении PDF: {str(e)}")
    
    # Provide more informative error message
    if not text.strip():
        # Check if document has images (might be scanned)
        has_images_in_doc = False
        try:
            with fitz.open(file_path) as doc:
                for page in doc:
                    if page.get_images():
                        has_images_in_doc = True
                        break
        except Exception:
            pass
        
        error_msg = (
            "Не удалось извлечь текст из PDF.\n\n"
            f"Проверено страниц: {total_pages}\n"
        )
        
        if has_images_in_doc:
            error_msg += (
                "⚠️ Обнаружены изображения в PDF.\n"
                "Возможно, это сканированный документ или текст встроен в картинки.\n\n"
            )
        
        error_msg += (
            "Возможные причины:\n"
            "• PDF содержит только изображения (текст встроен в картинки)\n"
            "• PDF является сканированной копией документа\n"
            "• Текст защищён от копирования\n"
            "• Файл повреждён или имеет нестандартный формат\n\n"
            "💡 Решения:\n"
            "• Используйте PDF с текстовым слоем (не сканированный)\n"
            "• Конвертируйте изображения в текст с помощью OCR-программ\n"
            "• Проверьте, что файл открывается в обычном PDF-ридере\n"
            "• Попробуйте экспортировать презентацию как PDF из PowerPoint/Google Slides"
        )
        raise ValueError(error_msg)
    
    # Warn if only few pages have text
    if pages_with_text < total_pages * 0.3 and total_pages > 1:
        print(f"[WARNING] Текст извлечён только с {pages_with_text} из {total_pages} страниц")
    
    return text.strip()


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file using python-pptx with improved extraction."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            
            # Extract from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shape_text = shape.text.strip()
                    # Skip very short text that might be decorative
                    if len(shape_text) > 2:
                        slide_text_parts.append(shape_text)
            
            # Also try to extract from tables
            for shape in slide.shapes:
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text_parts.append(" | ".join(row_text))
            
            if slide_text_parts:
                slide_text = "\n".join(slide_text_parts)
                text += f"\n--- Слайд {slide_num} ---\n"
                text += slide_text + "\n"
                
    except Exception as e:
        raise ValueError(f"Ошибка при чтении PPTX: {str(e)}")
    
    if not text.strip():
        raise ValueError("Не удалось извлечь текст из PPTX. Убедитесь, что файл содержит текстовую информацию.")
    
    return text.strip()
