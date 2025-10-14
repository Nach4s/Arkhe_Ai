import fitz  # PyMuPDF
from pptx import Presentation

def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from PDF or PPTX file.
    
    Args:
        file_path: Path to the file
        
    Returns:
        Extracted text as string
        
    Raises:
        ValueError: If file format is not supported
    """
    if file_path.lower().endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.lower().endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Поддерживаются только PDF и PPTX файлы.")


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file using PyMuPDF."""
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():
                    text += f"\n--- Страница {page_num} ---\n"
                    text += page_text
    except Exception as e:
        raise ValueError(f"Ошибка при чтении PDF: {str(e)}")
    
    return text.strip()


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file using python-pptx."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                text += f"\n--- Слайд {slide_num} ---\n"
                text += slide_text
    except Exception as e:
        raise ValueError(f"Ошибка при чтении PPTX: {str(e)}")
    
    return text.strip()
